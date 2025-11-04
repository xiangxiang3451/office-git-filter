import os
import subprocess
import tempfile
from .base_filter import BaseFilter
import sys

class OfficeFilter(BaseFilter):
    """Office文档过滤器"""

    def __init__(self):
        super().__init__()
        self.supported_formats = [
            '.doc', '.docx', '.xls', '.xlsx',
            '.ppt', '.pptx', '.odt', '.ods', '.odp'
        ]

    def to_text(self, file_path: str) -> str:
        """转换Office文档为文本"""
        ext = os.path.splitext(file_path)[1].lower()

        # 尝试使用LibreOffice
        text = self._use_libreoffice(file_path)
        if text:
            return text

        # 特定格式的备用方法
        if ext == '.doc':
            return self._use_antiword(file_path)
        elif ext == '.docx':
            return self._use_docx_library(file_path)
        elif ext in ['.xls', '.xlsx']:
            return self._use_excel_library(file_path)

        elif ext in ['.ppt', '.pptx']:
            return self._use_pptx_library(file_path)

        return f"[Office文件: {os.path.basename(file_path)}]"

    def _use_libreoffice(self, file_path: str) -> str:
        """使用LibreOffice进行转换"""
        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                # 转换为文本
                result = subprocess.run([
                    'soffice', '--headless', '--convert-to', 'txt:Text',
                    '--outdir', temp_dir, file_path
                ], capture_output=True, check=True)

                # 读取转换后的文件
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                txt_file = os.path.join(temp_dir, f"{base_name}.txt")

                if os.path.exists(txt_file):
                    with open(txt_file, 'r', encoding='utf-8', errors='ignore') as f:
                        return f.read().strip()

            return ""
        except (subprocess.CalledProcessError, FileNotFoundError):
            return ""

    def _use_antiword(self, file_path: str) -> str:
        """使用antiword处理.doc文件"""
        try:
            result = subprocess.run(
                ['antiword', file_path],
                capture_output=True,
                text=True,
                check=True
            )
            return result.stdout.strip()
        except (subprocess.CalledProcessError, FileNotFoundError):
            return self._use_catdoc(file_path)

    def _use_catdoc(self, file_path: str) -> str:
        """使用catdoc处理.doc文件"""
        try:
            result = subprocess.run(
                ['catdoc', file_path],
                capture_output=True,
                text=True,
                check=True
            )
            return result.stdout.strip()
        except (subprocess.CalledProcessError, FileNotFoundError):
            return ""

    def _use_docx_library(self, file_path: str) -> str:
        """使用python-docx处理.docx文件"""
        try:
            import docx
            doc = docx.Document(file_path)
            text = []
            for paragraph in doc.paragraphs:
                text.append(paragraph.text)
            return '\n'.join(text)
        except ImportError:
            return ""

    def _use_excel_library(self, file_path: str) -> str:
        """使用openpyxl处理Excel文件"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            text = []

            print(f"调试: 开始处理Excel文件 {file_path}", file=sys.stderr)

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text.append(f"=== 工作表: {sheet_name} ===")

                has_content = False
                for row_num, row in enumerate(sheet.iter_rows(values_only=True), 1):
                    # 过滤None值并转换为字符串
                    row_data = []
                    for cell in row:
                        if cell is not None:
                            cell_str = str(cell).strip()
                            if cell_str:  # 只添加非空字符串
                                row_data.append(cell_str)

                    if row_data:  # 只添加非空行
                        text.append(f"行{row_num}: {' | '.join(row_data)}")
                        has_content = True

                if not has_content:
                    text.append("(空工作表)")
                text.append("")  # 空行分隔

            result = '\n'.join(text)
            print(f"调试: Excel转换结果长度 {len(result)}", file=sys.stderr)
            return result if result.strip() else f"[Excel文件无文本内容: {os.path.basename(file_path)}]"

        except ImportError:
            return f"[未安装openpyxl库: {os.path.basename(file_path)}]"
        except Exception as e:
            return f"[处理Excel文件失败: {str(e)}]"


    def _use_pptx_library(self, file_path: str) -> str:
        """使用python-pptx处理PPTX文件"""
        try:
            import pptx
            presentation = pptx.Presentation(file_path)
            text = []

            print(f"调试: 开始处理PPTX文件 {file_path}", file=sys.stderr)

            for slide_num, slide in enumerate(presentation.slides, 1):
                text.append(f"=== 幻灯片 {slide_num} ===")

                has_content = False

                # 提取幻灯片标题
                if slide.shapes.title and slide.shapes.title.text.strip():
                    text.append(f"标题: {slide.shapes.title.text.strip()}")
                    has_content = True

                # 提取所有形状中的文本
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # 跳过标题（已经处理过）
                        if shape != slide.shapes.title:
                            text.append(f"内容: {shape.text.strip()}")
                            has_content = True

                if not has_content:
                    text.append("(空幻灯片)")
                text.append("")  # 空行分隔

            result = '\n'.join(text)
            print(f"调试: PPTX转换结果长度 {len(result)}", file=sys.stderr)
            return result if result.strip() else f"[PPTX文件无文本内容: {os.path.basename(file_path)}]"

        except ImportError:
            return f"[未安装python-pptx库: {os.path.basename(file_path)}]"
        except Exception as e:
            return f"[处理PPTX文件失败: {str(e)}]"